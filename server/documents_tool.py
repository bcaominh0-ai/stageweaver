"""
mcp_document_processing.py
FastMCP server exposing a Document‑processing tool
that works without the camel package.
"""

# --------------------------------------------------------------------------- #
#  Imports
# --------------------------------------------------------------------------- #
import asyncio, base64, os, io, json, subprocess, tempfile
from typing import Tuple, Optional, List

from loguru import logger
from retry import retry

from mcp.server.fastmcp import FastMCP
import anyio

# --- your own helper toolkits ------------------------------------------------ #
#   (provide these scripts in the Python path)
from image_tool import ask_question_about_image
from excel_tool import ExcelToolkit
# --- third‑party libs already used ------------------------------------------ #
import assemblyai as aai
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx import Presentation
from PIL import Image
from docx2markdown._docx_to_markdown import docx_to_markdown
import requests
import xmltodict
import nest_asyncio
nest_asyncio.apply()

from dotenv import load_dotenv
load_dotenv(".env")


def _env_nonempty(name: str) -> str:
    value = os.getenv(name)
    if value is None:
        return ""
    return value.strip()


AUDIO_TRANSCRIPTION_MODEL = _env_nonempty("AUDIO_TRANSCRIPTION_MODEL") or "whisper-1"
EXPLICIT_AUDIO_TRANSCRIPTION_OVERRIDE = any(
    _env_nonempty(name)
    for name in ("AUDIO_TRANSCRIPTION_API_KEY", "AUDIO_TRANSCRIPTION_BASE_URL", "AUDIO_TRANSCRIPTION_MODEL")
)
SOMARK_MODEL = _env_nonempty("SOMARK_MODEL") or "somark"
EXPLICIT_SOMARK_OVERRIDE = any(
    _env_nonempty(name)
    for name in ("SOMARK_API_KEY", "SOMARK_BASE_URL", "SOMARK_MODEL")
)
SOMARK_SUPPORTED_EXTENSIONS = {
    ".pdf", ".png", ".jpg", ".jpeg", ".bmp", ".tiff", ".jp2", ".dib",
    ".ppm", ".pgm", ".pbm", ".gif", ".heic", ".heif", ".webp", ".xpm",
    ".tga", ".dds", ".xbm",
}


# --------------------------------------------------------------------------- #
#  Toolkit implementation (no camel.BaseToolkit!)
# --------------------------------------------------------------------------- #
class DocumentProcessingToolkit:
    """
    This tool exposes a **general‑purpose document‑processing endpoint** that
    converts almost any common file you point it to into **clean, readable
    text or Markdown**.  It is useful whenever an agent needs to “look inside”
    an arbitrary file before reasoning over its contents.

    Conceptually, you can think of it as an *all‑in‑one* “open the file and
    give me the text” utility:

    • **Images (.jpg / .jpeg / .png)** – runs a vision model and returns a
    detailed caption.  
    • **Audio (.mp3 / .wav / .m4a)** – performs automatic transcription.  
    • **PowerPoint (.pptx)** – pulls every textbox, captions each embedded
    image, and preserves the slide order.  
    • **Spreadsheets (.xls / .xlsx / .csv)** – dumps cell values in a
    readable, row‑wise layout.  
    • **ZIP archives** – unpacks the archive and lists the extracted files.  
    • **Plain text‑like formats (.py / .txt)** – simply reads the file.  
    • **JSON, JSONL, JSON‑LD** – returns the parsed JSON structure.  
    • **XML** – converts to a Python dict (falls back to raw XML on error).  
    • **Word (.docx)** – converts the entire document to Markdown.  
    • **PDF or supported image formats** – attempts Somark extraction first,
      then a plain-text PDF fallback if Somark fails.

    Typical downstream tasks include:

    - Letting an LLM **summarise** or **answer questions about** a
    presentation, spreadsheet, contract, or research paper.  
    - Quickly **indexing** large document batches for semantic search.  
    - **Captioning image assets** to improve accessibility.  
    - Turning “opaque” binary files into human‑readable text for diffing or
    version control.

    ### Args
    `document_path` *(str)* – A fully‑qualified **local file path** that the
    server can access (e.g. `/home/user/input/report.pdf`).  Network URLs are
    *not* accepted.

    ### Returns
    `str` – On success, a **plain‑text or Markdown** representation of the
    file’s meaningful content.  
    If the file type is unsupported or an extraction error occurs, the tool
    raises an exception containing a diagnostic message.
    """

    def __init__(self, cache_dir: Optional[str] = None):
        self.excel_tool = ExcelToolkit()
        self.cache_dir = cache_dir or "tmp/"

    # --------------------------------------------------------------------- #
    #  Public façade
    # --------------------------------------------------------------------- #
    @retry(Exception,tries=5, delay=2, backoff=2)
    def extract_document_content(self, document_path: str) -> Tuple[bool, str]:
        logger.debug(f"[extract_document_content] {document_path=}")

        # 1. Images ----------------------------------------------------------------
        if document_path.lower().endswith((".jpg", ".jpeg", ".png")):
            caption = asyncio.run(
                ask_question_about_image(
                    document_path,
                    "Please make a detailed caption about the image."
                )
            )
            return True, caption

        # 2. Audio -----------------------------------------------------------------
        if document_path.lower().endswith((".mp3", ".wav", ".m4a")):
            return True, self._transcribe_audio(document_path)

        # 3. PPTX ------------------------------------------------------------------
        if document_path.lower().endswith(".pptx"):
            return True, asyncio.run(self._extract_pptx(document_path))

        # 4. Spreadsheets -----------------------------------------------------------
        if document_path.lower().endswith((".xls", ".xlsx", ".csv")):
            return True, self.excel_tool.extract_excel_content(document_path)

        # 5. Zip --------------------------------------------------------------------
        if document_path.lower().endswith(".zip"):
            return True, f"The extracted files are: {self._unzip_file(document_path)}"

        # 6. Simple text‑like formats ----------------------------------------------
        simple_readers = {
            ".py":  lambda p: open(p, encoding="utf‑8").read(),
            ".txt": lambda p: open(p, encoding="utf‑8").read(),
        }
        if any(document_path.lower().endswith(ext) for ext in simple_readers):
            reader = simple_readers[os.path.splitext(document_path)[1]]
            return True, reader(document_path)

        # 7. JSON                                                                   #
        if document_path.lower().endswith((".json", ".jsonl", ".jsonld")):
            return True, self._extract_json(document_path, encoding="utf‑8")
        

        # 8. XML                                                                    #
        if document_path.lower().endswith(".xml"):
            data = open(document_path, encoding="utf‑8").read()
            try:
                return True, xmltodict.parse(data)
            except Exception:
                return True, data

        # 9. DOCX → markdown -------------------------------------------------------
        if document_path.lower().endswith(".docx"):
            md_path = f"{os.path.basename(document_path)}.md"
            docx_to_markdown(document_path, md_path)
            return True, open(md_path, encoding="utf‑8").read()

        # 10. Fallback — optional Somark / PDF text ---------------------------------
        return self._try_somark_then_fallback(document_path)

    # ------------------------------------------------------------------------- #
    #  helpers
    # ------------------------------------------------------------------------- #
    def _extract_json(self, json_path: str, encoding: str = "utf‑8") -> str:
        with open(json_path, 'r', encoding=encoding) as f:
            if json_path.lower().endswith((".json",".jsonld")):
                return json.load(f)  
            elif json_path.lower().endswith(".jsonl"):
                return [json.loads(line) for line in f]                    

    async def _extract_pptx(self, pptx_path: str) -> str:
        prs = Presentation(pptx_path)
        base = pptx_path.rsplit(".", 1)[0]
        out = []

        for slide_idx, slide in enumerate(prs.slides, 1):
            txt = [f"Page {slide_idx}"]
            captions = []
            img_count = 0

            for shape_idx, shape in enumerate(slide.shapes):
                if getattr(shape, "text", "").strip():
                    txt.append(shape.text.strip())

                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    img_count += 1
                    img = Image.open(io.BytesIO(shape.image.blob))
                    img_path = f"{base}_slide_{slide_idx}_img_{shape_idx}.png"
                    img.save(img_path)
                    captions.append(
                        f"Image {img_count}: "
                        + await ask_question_about_image(
                            img_path, "Please make a detailed caption about the image."
                        )
                    )

            out.append("\n".join(txt + captions))

        return "\n\n".join(out)

    def _transcribe_audio(self, document_path: str) -> str:
        if EXPLICIT_AUDIO_TRANSCRIPTION_OVERRIDE:
            api_key = _env_nonempty("AUDIO_TRANSCRIPTION_API_KEY")
            if not api_key:
                raise RuntimeError(
                    "AUDIO_TRANSCRIPTION_* override is active but AUDIO_TRANSCRIPTION_API_KEY is missing. "
                    "Fill AUDIO_TRANSCRIPTION_API_KEY for the configured transcription endpoint."
                )
            base_url = _env_nonempty("AUDIO_TRANSCRIPTION_BASE_URL") or "https://www.dmxapi.cn/v1"
            endpoint = base_url.rstrip("/")
            if not endpoint.endswith("/responses"):
                endpoint = f"{endpoint}/responses"

            upload_path = document_path
            cleanup_path = None
            audio_format = os.path.splitext(document_path)[1].lstrip(".").lower() or "wav"

            # DMXAPI accepts wav/mp3 cleanly in practice. Convert m4a inputs to wav
            # before upload so documents_tool keeps supporting local .m4a files.
            if audio_format == "m4a":
                with tempfile.NamedTemporaryFile(suffix=".wav", delete=False) as tmp_file:
                    cleanup_path = tmp_file.name
                subprocess.run(
                    ["ffmpeg", "-y", "-i", document_path, cleanup_path],
                    check=True,
                    stdout=subprocess.DEVNULL,
                    stderr=subprocess.DEVNULL,
                )
                upload_path = cleanup_path
                audio_format = "wav"

            try:
                with open(upload_path, "rb") as audio_file:
                    audio_b64 = base64.b64encode(audio_file.read()).decode("utf-8")

                payload = {
                    "model": AUDIO_TRANSCRIPTION_MODEL,
                    "input": [
                        {
                            "role": "user",
                            "content": [
                                {
                                    "type": "input_audio",
                                    "input_audio": {
                                        "data": f"data:;base64,{audio_b64}",
                                        "format": audio_format,
                                    },
                                },
                                {
                                    "type": "text",
                                    "text": "请转写这段音频，尽量逐字输出原始内容，不要添加解释。",
                                },
                            ],
                        }
                    ],
                    "stream": True,
                    "stream_options": {"include_usage": True},
                    "modalities": ["text"],
                }
                headers = {
                    "Content-Type": "application/json",
                    "Authorization": api_key,
                }

                response = requests.post(
                    endpoint,
                    headers=headers,
                    json=payload,
                    stream=True,
                    timeout=120,
                )
                if not response.ok:
                    raise RuntimeError(
                        f"Audio transcription request failed with HTTP {response.status_code}: {response.text}"
                    )

                text_parts: List[str] = []
                current_event = None
                for line in response.iter_lines():
                    if not line:
                        continue
                    line_text = line.decode("utf-8").strip()
                    if line_text.startswith("event: "):
                        current_event = line_text[7:]
                        continue
                    if not line_text.startswith("data: "):
                        continue
                    data_str = line_text[6:]
                    if data_str == "[DONE]":
                        continue
                    try:
                        json_data = json.loads(data_str)
                    except json.JSONDecodeError:
                        continue

                    if current_event == "response.output_text.delta":
                        delta = json_data.get("delta", "")
                        if isinstance(delta, str) and delta:
                            text_parts.append(delta)
                        continue

                    if current_event == "response.output_item.done":
                        item = json_data.get("item", {})
                        for content in item.get("content", []):
                            if content.get("type") != "output_text":
                                continue
                            text = content.get("text", "")
                            if isinstance(text, str) and text:
                                text_parts.append(text)

                transcript = "".join(text_parts).strip()
                if transcript:
                    return transcript
                raise RuntimeError("Audio transcription completed but returned no text.")
            finally:
                if cleanup_path and os.path.exists(cleanup_path):
                    os.remove(cleanup_path)

        assembly_key = os.getenv("ASSEMBLYAI_API_KEY")
        if not assembly_key:
            raise RuntimeError(
                "Audio transcription is not configured. Set AUDIO_TRANSCRIPTION_* or ASSEMBLYAI_API_KEY."
            )
        aai.settings.api_key = assembly_key
        config = aai.TranscriptionConfig(speech_model=aai.SpeechModel.best)
        transcript = aai.Transcriber(config=config).transcribe(document_path)
        logger.info(transcript.text)
        if transcript.status == "error":
            raise RuntimeError(f"Transcription failed: {transcript.error}")
        return transcript.text

    @staticmethod
    def _extract_pdf_with_pypdf2(path: str) -> str:
        from PyPDF2 import PdfReader

        return "\n".join(
            page.extract_text() or ""
            for page in PdfReader(open(path, "rb")).pages
        )

    def _try_somark_then_fallback(self, path: str) -> Tuple[bool, str]:
        ext = os.path.splitext(path)[1].lower()
        try:
            if ext not in SOMARK_SUPPORTED_EXTENSIONS:
                raise RuntimeError(f"Unsupported file type for Somark fallback: {ext or 'unknown'}")
            if not _env_nonempty("SOMARK_API_KEY"):
                raise RuntimeError("Somark fallback is not configured. Set SOMARK_API_KEY.")
            text = self._extract_with_somark(path, output_format="markdown")
            return True, text
        except Exception as e:
            logger.warning(f"Somark failed: {e}")
            if path.lower().endswith(".pdf"):
                try:
                    text = self._extract_pdf_with_pypdf2(path)
                    if text.strip():
                        return True, text
                    return False, f"PDF fallback produced no text after Somark failed: {e}"
                except Exception as e2:
                    return False, f"PDF fallback failed after Somark error ({e}): {e2}"
            return False, f"Unsupported file type or Somark processing error: {e}"

    def _extract_with_somark(self, path: str, output_format: str = "markdown") -> str:
        api_key = _env_nonempty("SOMARK_API_KEY")
        if not api_key:
            raise RuntimeError("Missing SOMARK_API_KEY.")

        base_url = _env_nonempty("SOMARK_BASE_URL") or "https://www.dmxapi.cn/v1/responses"
        endpoint = base_url.rstrip("/")
        if not endpoint.endswith("/responses"):
            endpoint = f"{endpoint}/responses"

        with open(path, "rb") as f:
            file_base64 = base64.b64encode(f.read()).decode("utf-8")

        headers = {
            "Authorization": api_key,
            "Content-Type": "application/json",
        }
        payload = {
            "model": SOMARK_MODEL,
            "input": "json",
            "file": file_base64,
        }

        response = requests.post(endpoint, headers=headers, json=payload, timeout=180)
        if not response.ok:
            raise RuntimeError(f"Somark request failed with HTTP {response.status_code}: {response.text}")

        data = response.json()
        if data.get("code") != 0:
            raise RuntimeError(f"Somark returned code={data.get('code')}: {data.get('message') or data}")

        result = ((data.get("data") or {}).get("result") or {})
        outputs = result.get("outputs") or {}
        json_output = outputs.get("json") or {}
        if output_format == "json":
            return json.dumps(json_output, ensure_ascii=False, indent=2)
        rendered = self._render_somark_markdown(json_output)
        if rendered.strip():
            return rendered
        raise RuntimeError("Somark returned no usable text blocks.")

    @staticmethod
    def _render_somark_markdown(parsed: dict) -> str:
        page_texts: List[str] = []
        for page in parsed.get("pages", []):
            blocks = page.get("blocks", [])
            block_texts: List[str] = []
            for block in blocks:
                if block.get("display") is False:
                    continue
                content = block.get("content", "")
                if isinstance(content, str) and content.strip():
                    block_texts.append(content.strip())
            page_text = "\n".join(block_texts).strip()
            if page_text:
                page_num = page.get("page_num")
                if page_num is not None:
                    page_texts.append(f"Page {page_num + 1}\n{page_text}")
                else:
                    page_texts.append(page_text)
        return "\n\n".join(page_texts)

    def _unzip_file(self, zip_path: str) -> List[str]:
        dest = os.path.join(self.cache_dir, os.path.splitext(os.path.basename(zip_path))[0])
        os.makedirs(dest, exist_ok=True)
        subprocess.run(["unzip", "-o", zip_path, "-d", dest], check=True)
        return [os.path.join(r, f) for r, _, fs in os.walk(dest) for f in fs]


# --------------------------------------------------------------------------- #
#  FastMCP server
# --------------------------------------------------------------------------- #
mcp = FastMCP("document_processing")
toolkit = DocumentProcessingToolkit()


@mcp.tool()
async def process_document(document_path: str) -> str:
    """
    Process a document at the given *document_path*. The document can be multimedia
    (image, audio), a presentation (PPTX), a spreadsheet, a ZIP archive,
    a text file, JSON, XML, Word document, or PDF.
   Return the extracted text / markdown representation of *document_path*.
   supported formats include .png, .jpeg, .jpg, .mp3, .m4a, .PPTX, .xlsx, .csv, .txt, .json, .jsonl, .jsonld, .zip, .xml, .docx, and .pdf.
    """
    success, content = await anyio.to_thread.run_sync(
        toolkit.extract_document_content, document_path
    )
    if not success:
        raise ValueError(content)
    return content


# --------------------------------------------------------------------------- #
#  Entrypoint
# --------------------------------------------------------------------------- #
if __name__ == "__main__":
    mcp.run(transport="stdio")
